"""Tool definitions and dispatch for LLM function calling."""

from typing import Any

from memory import write_memory, read_memory, list_memories
from sandbox import run_code
from search import web_search, format_search_results

# OpenAI function-calling format tool definitions
TOOL_DEFINITIONS = [
    {
        "type": "function",
        "function": {
            "name": "web_search",
            "description": "Search the web for current information. Use when the user asks about recent events, news, or anything that may not be in your training data.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "The search query",
                    }
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_memory",
            "description": "Read a previously saved memory file. Use when the user asks about something you may have remembered.",
            "parameters": {
                "type": "object",
                "properties": {
                    "filename": {
                        "type": "string",
                        "description": "The memory filename to read (e.g. 'user_name.txt')",
                    },
                    "subdir": {
                        "type": "string",
                        "enum": ["facts", "conversations", "notes"],
                        "description": "Memory subdirectory (default: facts)",
                    },
                },
                "required": ["filename"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_memory",
            "description": "Save information to persistent memory. ONLY use when the user explicitly asks you to remember or save something. Do not proactively save observations about the conversation.",
            "parameters": {
                "type": "object",
                "properties": {
                    "filename": {
                        "type": "string",
                        "description": "Filename for the memory (e.g. 'user_name.txt')",
                    },
                    "content": {
                        "type": "string",
                        "description": "The content to save",
                    },
                    "subdir": {
                        "type": "string",
                        "enum": ["facts", "conversations", "notes"],
                        "description": "Memory subdirectory (default: facts)",
                    },
                },
                "required": ["filename", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "list_memories",
            "description": "List all saved memory files. Use when the user asks what you remember.",
            "parameters": {
                "type": "object",
                "properties": {
                    "subdir": {
                        "type": "string",
                        "enum": ["facts", "conversations", "notes"],
                        "description": "Filter to specific subdirectory (optional)",
                    },
                },
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "run_code",
            "description": "Execute code in a sandboxed container. Use ONLY when the user asks you to run code, do a calculation, or when a task clearly requires computation. Do NOT use this for document generation — use generate_document instead. Supported languages: python (with numpy, pandas, matplotlib, sympy, scipy), javascript (Node.js), bash, c, cpp, go, lua. No network access.",
            "parameters": {
                "type": "object",
                "properties": {
                    "code": {
                        "type": "string",
                        "description": "Source code to execute",
                    },
                    "language": {
                        "type": "string",
                        "description": "Programming language (default: python)",
                        "enum": ["python", "javascript", "bash", "c", "cpp", "go", "lua"],
                    },
                    "timeout": {
                        "type": "integer",
                        "description": "Execution timeout in seconds (default 10, max 30). Use higher values for compiled languages.",
                        "minimum": 1,
                        "maximum": 30,
                    },
                },
                "required": ["code"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "generate_document",
            "description": "Generate a downloadable document file. Use this when the user asks to create a PDF, Word doc, Excel spreadsheet, PowerPoint, or CSV. Provide the content as structured text — the system handles formatting and file creation. ALWAYS use this tool for document generation — never fabricate download links.",
            "parameters": {
                "type": "object",
                "properties": {
                    "format": {
                        "type": "string",
                        "description": "Output format",
                        "enum": ["pdf", "docx", "xlsx", "csv", "pptx", "txt"],
                    },
                    "title": {
                        "type": "string",
                        "description": "Document title",
                    },
                    "content": {
                        "type": "string",
                        "description": "Document content. For PDF/DOCX/TXT: plain text with paragraphs separated by newlines. For CSV/XLSX: rows separated by newlines, columns separated by | (first row is headers). For PPTX: slides separated by --- with first line as slide title.",
                    },
                },
                "required": ["format", "title", "content"],
            },
        },
    },
]


def _build_doc_code(fmt: str, title: str, content: str) -> str:
    """Build Python code to generate a document from structured content."""
    # Escape for safe injection into Python string
    t = title.replace("\\", "\\\\").replace("'", "\\'")
    c = content.replace("\\", "\\\\").replace("'", "\\'")

    if fmt == "txt":
        return f"""import os
os.makedirs('/tmp/output', exist_ok=True)
with open('/tmp/output/{t}.txt', 'w') as f:
    f.write('''{c}''')
print('Document created')
"""

    if fmt == "csv":
        return f"""import csv, os
os.makedirs('/tmp/output', exist_ok=True)
lines = '''{c}'''.strip().split('\\n')
with open('/tmp/output/{t}.csv', 'w', newline='') as f:
    writer = csv.writer(f)
    for line in lines:
        writer.writerow([col.strip() for col in line.split('|')])
print('CSV created')
"""

    if fmt == "pdf":
        return f"""import os
os.makedirs('/tmp/output', exist_ok=True)
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
doc = SimpleDocTemplate('/tmp/output/{t}.pdf', pagesize=letter)
styles = getSampleStyleSheet()
story = [Paragraph('''{t}''', styles['Title']), Spacer(1, 12)]
for para in '''{c}'''.split('\\n'):
    para = para.strip()
    if para:
        story.append(Paragraph(para, styles['BodyText']))
        story.append(Spacer(1, 6))
doc.build(story)
print('PDF created')
"""

    if fmt == "docx":
        return f"""import os
os.makedirs('/tmp/output', exist_ok=True)
from docx import Document
doc = Document()
doc.add_heading('''{t}''', level=1)
for para in '''{c}'''.split('\\n'):
    para = para.strip()
    if para:
        doc.add_paragraph(para)
doc.save('/tmp/output/{t}.docx')
print('DOCX created')
"""

    if fmt == "xlsx":
        return f"""import os
os.makedirs('/tmp/output', exist_ok=True)
from openpyxl import Workbook
from openpyxl.styles import Font
wb = Workbook()
ws = wb.active
ws.title = '''{t}'''
lines = '''{c}'''.strip().split('\\n')
for i, line in enumerate(lines, 1):
    cols = [col.strip() for col in line.split('|')]
    for j, val in enumerate(cols, 1):
        cell = ws.cell(row=i, column=j, value=val)
        if i == 1:
            cell.font = Font(bold=True)
wb.save('/tmp/output/{t}.xlsx')
print('XLSX created')
"""

    if fmt == "pptx":
        return f"""import os
os.makedirs('/tmp/output', exist_ok=True)
from pptx import Presentation
from pptx.util import Inches, Pt
prs = Presentation()
slides_raw = '''{c}'''.split('---')
for slide_text in slides_raw:
    lines = [l.strip() for l in slide_text.strip().split('\\n') if l.strip()]
    if not lines:
        continue
    slide_title = lines[0]
    body = '\\n'.join(lines[1:]) if len(lines) > 1 else ''
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = slide_title
    if body and slide.placeholders[1]:
        slide.placeholders[1].text = body
prs.save('/tmp/output/{t}.pptx')
print('PPTX created')
"""

    return "print('Unsupported format')"


async def execute_tool(name: str, arguments: dict[str, Any]) -> str:
    """Execute a tool by name and return the result as a string."""
    if name == "web_search":
        results = await web_search(arguments["query"])
        return format_search_results(results)

    elif name == "read_memory":
        return read_memory(
            arguments["filename"],
            arguments.get("subdir", "facts"),
        )

    elif name == "write_memory":
        return write_memory(
            arguments["filename"],
            arguments["content"],
            arguments.get("subdir", "facts"),
        )

    elif name == "list_memories":
        memories = list_memories(arguments.get("subdir"))
        if not memories:
            return "No memories saved yet."
        lines = ["Saved memories:\n"]
        for m in memories:
            lines.append(f"- {m['subdir']}/{m['filename']} ({m['size']} bytes)")
        return "\n".join(lines)

    elif name == "generate_document":
        fmt = arguments.get("format", "txt")
        title = arguments.get("title", "document")
        content = arguments.get("content", "")
        code = _build_doc_code(fmt, title, content)
        result = await run_code(code, "python", 20)
        parts = []
        if result["stderr"]:
            parts.append(f"Error: {result['stderr']}")
        if result.get("output_files"):
            for f in result["output_files"]:
                parts.append(f"Generated: [{f['filename']}]({f['url']})")
        if not parts:
            parts.append("Document generation failed — no output file produced.")
        return "\n".join(parts)

    elif name == "run_code":
        result = await run_code(
            arguments["code"],
            arguments.get("language", "python"),
            arguments.get("timeout", 10),
        )
        parts = []
        if result["timed_out"]:
            parts.append(f"[TIMED OUT after {arguments.get('timeout', 10)}s]")
        if result["stdout"]:
            parts.append(f"stdout:\n{result['stdout']}")
        if result["stderr"]:
            parts.append(f"stderr:\n{result['stderr']}")
        if result.get("output_files"):
            file_lines = ["Generated files:"]
            for f in result["output_files"]:
                file_lines.append(f"- {f['filename']}: {f['url']}")
            parts.append("\n".join(file_lines))
        if not parts:
            parts.append(f"(no output, exit code {result['exit_code']})")
        return "\n".join(parts)

    else:
        return f"Unknown tool: {name}"
